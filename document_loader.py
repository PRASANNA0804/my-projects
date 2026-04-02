"""
document_loader.py — Universal Multi-Format Document Loader
============================================================
Supported formats:
  .txt  — Plain text (auto-detects encoding)
  .md   — Markdown (strips syntax, keeps content)
  .pdf  — PDF via PyMuPDF (handles scanned + text PDFs)
  .docx — Word documents via python-docx
  .xlsx — Excel spreadsheets via openpyxl
  .xls  — Legacy Excel (basic support via openpyxl)
  .pptx — PowerPoint presentations via python-pptx
  .csv  — CSV files (tabular → readable text)
  .json — JSON files (pretty-printed key traversal)
  .rtf  — Rich Text Format via striprtf

Each loader returns raw text that is then chunked by ingest.py.
"""

from __future__ import annotations
import base64
import csv
import json
import io
import os
from pathlib import Path
from typing import Any, Callable

# ── Format registry (populated by decorators below) ──────────────────────────
_LOADERS: dict[str, Callable[[Path], str]] = {}

def loader(*extensions: str):
    """Decorator that registers a function as the loader for given extensions."""
    def decorator(fn: Callable[[Path], str]):
        for ext in extensions:
            _LOADERS[ext.lower()] = fn
        return fn
    return decorator


# ── .txt / .md / .log — Plain text ───────────────────────────────────────────
@loader(".txt", ".log", ".text")
def load_text(path: Path) -> str:
    try:
        import chardet
        raw   = path.read_bytes()
        enc   = str(chardet.detect(raw).get("encoding") or "utf-8")
        return raw.decode(enc, errors="replace")
    except ImportError:
        return path.read_text(encoding="utf-8", errors="replace")


@loader(".md", ".markdown")
def load_markdown(path: Path) -> str:
    """Strip markdown syntax; keep readable content."""
    text = path.read_text(encoding="utf-8", errors="replace")
    # Remove code fences, images, links — keep link text
    import re
    text = re.sub(r"```[\s\S]*?```", "", text)           # code blocks
    text = re.sub(r"`[^`]+`", lambda m: m.group()[1:-1], text)  # inline code
    text = re.sub(r"!\[.*?\]\(.*?\)", "", text)           # images
    text = re.sub(r"\[([^\]]+)\]\([^\)]+\)", r"\1", text) # links → text
    text = re.sub(r"^#{1,6}\s+", "", text, flags=re.MULTILINE)  # headings
    text = re.sub(r"\*{1,2}([^*]+)\*{1,2}", r"\1", text)  # bold/italic
    return text.strip()


# ── .pdf — PyMuPDF ────────────────────────────────────────────────────────────
@loader(".pdf")
def load_pdf(path: Path) -> str:
    try:
        import fitz  # PyMuPDF
        doc   = fitz.open(str(path))
        pages = []
        for i in range(len(doc)):
            page = doc.load_page(i)
            text: str = str(page.get_text("text"))
            if text.strip():
                pages.append(f"[Page {i+1}]\n{text}")
        doc.close()
        return "\n\n".join(pages) if pages else ""
    except ImportError:
        raise RuntimeError("Install PyMuPDF: pip install pymupdf")


# ── .docx — python-docx ───────────────────────────────────────────────────────
@loader(".docx")
def load_docx(path: Path) -> str:
    try:
        from docx import Document
        doc    = Document(str(path))
        parts  = []

        # Body paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                style = (para.style.name or "") if para.style else ""
                prefix = ""
                if style.startswith("Heading"):
                    level  = style.replace("Heading ", "").strip()
                    prefix = "#" * int(level) + " " if level.isdigit() else ""
                parts.append(f"{prefix}{para.text.strip()}")

        # Tables
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                rows.append(" | ".join(cells))
            if rows:
                parts.append("\n".join(rows))

        return "\n\n".join(parts)
    except ImportError:
        raise RuntimeError("Install python-docx: pip install python-docx")


# ── .xlsx / .xls — openpyxl ──────────────────────────────────────────────────
@loader(".xlsx", ".xls")
def load_excel(path: Path) -> str:
    try:
        import openpyxl
        wb    = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws: Any = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):
                    rows.append(" | ".join(cells))
            if rows:
                parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(parts)
    except ImportError:
        raise RuntimeError("Install openpyxl: pip install openpyxl")


# ── .pptx — python-pptx ──────────────────────────────────────────────────────
def _ocr_image_bytes(image_bytes: bytes) -> str:
    """
    Extract text from raw image bytes using pytesseract OCR.
    Returns empty string silently if pytesseract / Tesseract is not installed.
    """
    try:
        import pytesseract
        from PIL import Image
        # Point to Tesseract binary explicitly (Windows default install path)
        _TESSERACT_PATHS = [
            r"C:\Program Files\Tesseract-OCR\tesseract.exe",
            r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
        ]
        for _p in _TESSERACT_PATHS:
            if os.path.isfile(_p):
                pytesseract.pytesseract.tesseract_cmd = _p
                break
        img = Image.open(io.BytesIO(image_bytes))
        return pytesseract.image_to_string(img).strip()
    except Exception:
        return ""


@loader(".pptx", ".ppt")
def load_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs   = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_parts: list[str] = []

            for shape in slide.shapes:
                # Regular text shapes
                shape_text = getattr(shape, "text", None)
                if shape_text and str(shape_text).strip():
                    slide_parts.append(str(shape_text).strip())

                # Picture shapes — attempt OCR
                # Use getattr so linked/external images don't raise AttributeError
                shape_img = getattr(shape, "image", None)
                if shape_img is not None:
                    try:
                        ocr_text = _ocr_image_bytes(shape_img.blob)
                        if ocr_text:
                            slide_parts.append(f"[Image text: {ocr_text}]")
                    except Exception:
                        pass

            if slide_parts:
                parts.append(f"[Slide {i}]\n" + "\n".join(slide_parts))
        return "\n\n".join(parts)
    except ImportError:
        raise RuntimeError("Install python-pptx: pip install python-pptx")


# ── Image formats — Azure OpenAI Vision + pytesseract fallback ───────────────

# MIME type map for supported image extensions
_IMAGE_MIME: dict[str, str] = {
    ".jpg" : "image/jpeg",
    ".jpeg": "image/jpeg",
    ".png" : "image/png",
    ".gif" : "image/gif",
    ".bmp" : "image/bmp",
    ".webp": "image/webp",
}

# Vision prompt — tuned for architecture diagrams, screenshots, charts
_VISION_PROMPT = (
    "You are analyzing an image to extract ALL information for a searchable knowledge base.\n\n"
    "Describe everything you see in full detail:\n"
    "• Architecture / system diagram → list every component, service, database, queue, "
    "API, connection, data flow, protocol, and label. Explain the overall structure, "
    "relationships between components, and how data moves through the system.\n"
    "• Screenshot → transcribe every visible text, heading, field, value, button, "
    "menu item, error message, and UI element exactly as shown.\n"
    "• Chart / graph → describe the title, axes labels, legend, data series, "
    "specific values, and the key trends or conclusions.\n"
    "• Table → reproduce the full table content with all rows and column headers.\n"
    "• Photo / illustration → describe all visible content, text, annotations, "
    "and relevant context.\n\n"
    "Rules:\n"
    "- Include every text label, heading, and annotation visible in the image.\n"
    "- Be thorough and precise — your description is the ONLY way this image's "
    "content will be searchable in the knowledge base.\n"
    "- Do NOT say 'I can see' or add commentary — output structured descriptive text only."
)


def _describe_image_with_vision(image_bytes: bytes, mime_type: str) -> str:
    """
    Send image to Azure OpenAI (gpt-4o-mini supports vision) and return
    a rich text description. Reads credentials from environment / .env.
    """
    from dotenv import load_dotenv
    from openai import AzureOpenAI
    load_dotenv()

    endpoint   = os.getenv("AZURE_OPENAI_ENDPOINT", "")
    api_key    = os.getenv("AZURE_OPENAI_API_KEY", "")
    api_ver    = os.getenv("AZURE_OPENAI_API_VERSION", "2024-12-01-preview")
    deployment = os.getenv("AZURE_OPENAI_DEPLOYMENT", "gpt-4o-mini")

    if not endpoint or not api_key:
        raise RuntimeError("AZURE_OPENAI_ENDPOINT / AZURE_OPENAI_API_KEY not set in .env")

    b64_data = base64.b64encode(image_bytes).decode("utf-8")
    client   = AzureOpenAI(azure_endpoint=endpoint, api_key=api_key, api_version=api_ver)

    resp = client.chat.completions.create(
        model      = deployment,
        max_tokens = 1500,
        messages   = [{
            "role": "user",
            "content": [
                {
                    "type"     : "image_url",
                    "image_url": {
                        "url"   : f"data:{mime_type};base64,{b64_data}",
                        "detail": "high",
                    },
                },
                {"type": "text", "text": _VISION_PROMPT},
            ],
        }],
    )
    return (resp.choices[0].message.content or "").strip()


@loader(".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp")
def load_image(path: Path) -> str:
    """
    Extract content from an image file.

    Strategy (in order):
      1. Azure OpenAI vision  — understands diagrams, charts, screenshots,
                                spatial relationships, and component connections.
      2. pytesseract OCR      — fallback for plain text/screenshot images when
                                the vision API is unavailable.
      3. RuntimeError         — both paths failed; user sees a clear message.
    """
    image_bytes = path.read_bytes()
    mime_type   = _IMAGE_MIME.get(path.suffix.lower(), "image/png")

    # ── Primary: Azure OpenAI vision ──────────────────────────────────────────
    try:
        description = _describe_image_with_vision(image_bytes, mime_type)
        if description:
            return f"[Image: {path.name}]\n{description}"
    except Exception as vision_err:
        print(f"   ⚠  Vision API unavailable for '{path.name}': {vision_err}")
        print(     "      Falling back to OCR …")

    # ── Fallback: pytesseract OCR ─────────────────────────────────────────────
    ocr_text = _ocr_image_bytes(image_bytes)
    if ocr_text:
        return f"[Image OCR: {path.name}]\n{ocr_text}"

    raise RuntimeError(
        f"Could not extract content from '{path.name}'. "
        "Ensure Azure OpenAI is configured in .env, or install Tesseract OCR."
    )


# ── .csv ──────────────────────────────────────────────────────────────────────
@loader(".csv")
def load_csv(path: Path) -> str:
    text = path.read_text(encoding="utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    rows   = list(reader)
    if not rows:
        return ""
    header = rows[0]
    parts  = []
    for row in rows[1:]:
        if any(cell.strip() for cell in row):
            pairs = [f"{h}: {v}" for h, v in zip(header, row) if h.strip()]
            parts.append(", ".join(pairs))
    return f"[CSV: {len(rows)-1} rows, columns: {', '.join(header)}]\n\n" + "\n".join(parts)


# ── .json ─────────────────────────────────────────────────────────────────────
@loader(".json", ".jsonl")
def load_json(path: Path) -> str:
    text = path.read_text(encoding="utf-8", errors="replace")
    try:
        if path.suffix == ".jsonl":
            lines  = [json.loads(l) for l in text.splitlines() if l.strip()]
            return "\n\n".join(json.dumps(l, indent=2, ensure_ascii=False) for l in lines)
        obj = json.loads(text)
        return json.dumps(obj, indent=2, ensure_ascii=False)
    except json.JSONDecodeError:
        return text   # fall back to raw text


# ── .rtf — striprtf ──────────────────────────────────────────────────────────
@loader(".rtf")
def load_rtf(path: Path) -> str:
    try:
        from striprtf.striprtf import rtf_to_text
        return str(rtf_to_text(path.read_text(encoding="utf-8", errors="replace")) or "")
    except ImportError:
        raise RuntimeError("Install striprtf: pip install striprtf")


# ── Public API ────────────────────────────────────────────────────────────────
SUPPORTED_EXTENSIONS = sorted(_LOADERS.keys())

def is_supported(path: str | Path) -> bool:
    return Path(path).suffix.lower() in _LOADERS

def extract_text(path: str | Path) -> str:
    """
    Extract plain text from *path* using the appropriate loader.
    Raises ValueError for unsupported formats.
    """
    p   = Path(path)
    ext = p.suffix.lower()
    if ext not in _LOADERS:
        supported = ", ".join(SUPPORTED_EXTENSIONS)
        raise ValueError(f"Unsupported format '{ext}'. Supported: {supported}")
    return _LOADERS[ext](p)


if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        print(f"Supported formats: {', '.join(SUPPORTED_EXTENSIONS)}")
    else:
        text = extract_text(sys.argv[1])
        print(text[:2000])
